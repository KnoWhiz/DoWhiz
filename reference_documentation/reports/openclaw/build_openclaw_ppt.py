from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT_DIR = Path("reports/openclaw")
OUT_FILE = OUT_DIR / "openclaw_核心创新技术架构_详解_v1.pptx"


@dataclass(frozen=True)
class Palette:
    dark: str = "0B132B"
    navy: str = "1C2541"
    blue: str = "3A506B"
    teal: str = "5BC0BE"
    aqua: str = "2EC4B6"
    amber: str = "F4A261"
    sand: str = "F7F3E8"
    light: str = "F4F7FB"
    white: str = "FFFFFF"
    ink: str = "152238"
    gray: str = "5B677A"
    mint: str = "DDF5F2"
    rose: str = "FDEBE6"


P = Palette()


# 16:9 default size ~ 13.333 x 7.5 in
W = 13.333
H = 7.5


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip().replace("#", "")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_bg(slide, color_hex: str, with_deco: bool = True) -> None:
    bg = slide.shapes.add_shape(SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(color_hex)
    bg.line.fill.background()
    if not with_deco:
        return

    d1 = slide.shapes.add_shape(SHAPE.OVAL, Inches(W - 3.8), Inches(-0.9), Inches(4.6), Inches(4.6))
    d1.fill.solid()
    d1.fill.fore_color.rgb = rgb(P.teal)
    d1.fill.transparency = 80
    d1.line.fill.background()

    d2 = slide.shapes.add_shape(SHAPE.OVAL, Inches(-1.2), Inches(H - 2.7), Inches(3.6), Inches(3.6))
    d2.fill.solid()
    d2.fill.fore_color.rgb = rgb(P.amber)
    d2.fill.transparency = 84
    d2.line.fill.background()


def add_title(slide, text: str, subtitle: str | None = None, dark_mode: bool = False) -> None:
    text_len = len(text)
    if text_len >= 34:
        title_size = Pt(31)
        title_h = 1.65
    elif text_len >= 28:
        title_size = Pt(34)
        title_h = 1.6
    elif text_len >= 22:
        title_size = Pt(37)
        title_h = 1.65
    else:
        title_size = Pt(40)
        title_h = 1.4

    title_box = slide.shapes.add_textbox(Inches(0.72), Inches(0.58), Inches(12.05), Inches(title_h))
    tf = title_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = text
    p0.font.name = "PingFang SC"
    p0.font.bold = True
    p0.font.size = title_size
    p0.font.color.rgb = rgb(P.white if dark_mode else P.ink)

    # Long titles usually occupy more vertical space; hide subtitle to avoid content crowding.
    if subtitle and text_len < 30:
        sub_y = 0.58 + title_h + 0.05
        sub_box = slide.shapes.add_textbox(Inches(0.78), Inches(sub_y), Inches(12.1), Inches(0.95))
        stf = sub_box.text_frame
        stf.clear()
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "PingFang SC"
        sp.font.size = Pt(15)
        sp.font.color.rgb = rgb("D7E7F7" if dark_mode else "5F6F85")


def add_footer(slide, text: str, dark_mode: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(7.03), Inches(11.9), Inches(0.32))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "PingFang SC"
    p.font.size = Pt(10)
    p.font.color.rgb = rgb("D8E4F2" if dark_mode else "6F7F95")


def add_card(slide, x: float, y: float, w: float, h: float, title: str, body: Iterable[str], theme: str = "light") -> None:
    rect = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    if theme == "light":
        rect.fill.fore_color.rgb = rgb(P.white)
        line_color = P.light
        title_color = P.ink
        body_color = "334A64"
    elif theme == "mint":
        rect.fill.fore_color.rgb = rgb(P.mint)
        line_color = "C5ECE7"
        title_color = P.ink
        body_color = "2D425C"
    elif theme == "rose":
        rect.fill.fore_color.rgb = rgb(P.rose)
        line_color = "F9DCD2"
        title_color = P.ink
        body_color = "2D425C"
    else:
        rect.fill.fore_color.rgb = rgb("13253D")
        line_color = "24496E"
        title_color = P.white
        body_color = "D6E2F1"

    rect.line.color.rgb = rgb(line_color)
    rect.line.width = Pt(1.25)

    tbox = slide.shapes.add_textbox(Inches(x + 0.28), Inches(y + 0.18), Inches(w - 0.5), Inches(0.45))
    tf = tbox.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.name = "PingFang SC"
    p0.font.bold = True
    p0.font.size = Pt(18)
    p0.font.color.rgb = rgb(title_color)

    b = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.72), Inches(w - 0.55), Inches(h - 0.9))
    btf = b.text_frame
    btf.clear()
    for i, line in enumerate(body):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = f"• {line}"
        p.font.name = "PingFang SC"
        p.font.size = Pt(14.2)
        p.font.color.rgb = rgb(body_color)
        p.level = 0


def add_big_metric(slide, x: float, y: float, value: str, label: str, color_hex: str) -> None:
    box = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.0), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = rgb(P.white)
    box.line.color.rgb = rgb("DFE6F2")

    n = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.34), Inches(2.6), Inches(0.9))
    tf = n.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = value
    p.font.name = "Georgia"
    p.font.bold = True
    p.font.size = Pt(42)
    p.font.color.rgb = rgb(color_hex)
    p.alignment = PP_ALIGN.LEFT

    l = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 1.26), Inches(2.6), Inches(0.5))
    ltf = l.text_frame
    ltf.clear()
    lp = ltf.paragraphs[0]
    lp.text = label
    lp.font.name = "PingFang SC"
    lp.font.size = Pt(13)
    lp.font.color.rgb = rgb(P.gray)


def add_tier_row(slide, y: float, idx: int, text: str, width: float = 7.8) -> None:
    chip = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(y), Inches(0.75), Inches(0.45))
    chip.fill.solid()
    chip.fill.fore_color.rgb = rgb(P.dark)
    chip.line.fill.background()
    ctf = chip.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.text = str(idx)
    cp.alignment = PP_ALIGN.CENTER
    cp.font.name = "Georgia"
    cp.font.bold = True
    cp.font.size = Pt(14)
    cp.font.color.rgb = rgb(P.white)

    row = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(y), Inches(width), Inches(0.45))
    row.fill.solid()
    row.fill.fore_color.rgb = rgb(P.white)
    row.line.color.rgb = rgb("DFE6F2")
    rtf = row.text_frame
    rtf.clear()
    rp = rtf.paragraphs[0]
    rp.text = text
    rp.font.name = "PingFang SC"
    rp.font.size = Pt(13)
    rp.font.color.rgb = rgb(P.ink)


def build_ppt() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # Slide 1: Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, P.dark, with_deco=True)
    add_title(
        s,
        "OpenClaw 核心创新技术架构",
        "基于 external/openclaw 源码与官方文档的系统级深度拆解",
        dark_mode=True,
    )
    ribbon = s.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(3.35), Inches(8.3), Inches(1.95))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = rgb("112A46")
    ribbon.line.color.rgb = rgb("2A4A6E")

    rtf = ribbon.text_frame
    rtf.clear()
    p = rtf.paragraphs[0]
    p.text = "研究结论先行"
    p.font.name = "PingFang SC"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = rgb(P.aqua)

    for line in [
        "1) OpenClaw 的本质是“可治理 AI 网关”，不是单通道 Bot。",
        "2) 核心壁垒在于：统一控制平面 + 确定性路由 + 分层安全模型。",
        "3) 对 DoWhiz 的高价值迁移点：会话键规范、路由层次、运行时策略引擎。",
    ]:
        pp = rtf.add_paragraph()
        pp.text = line
        pp.font.name = "PingFang SC"
        pp.font.size = Pt(14)
        pp.font.color.rgb = rgb("D6E2F1")

    tag = s.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(10.1), Inches(5.8), Inches(2.35), Inches(1.0))
    tag.fill.solid()
    tag.fill.fore_color.rgb = rgb(P.amber)
    tag.line.fill.background()
    ttf = tag.text_frame
    ttf.clear()
    tp = ttf.paragraphs[0]
    tp.text = "18 页\n架构深拆"
    tp.alignment = PP_ALIGN.CENTER
    tp.font.name = "PingFang SC"
    tp.font.bold = True
    tp.font.size = Pt(20)
    tp.font.color.rgb = rgb(P.dark)

    add_footer(s, "来源：OpenClaw 源码仓库（v2026.2.23）+ docs.openclaw.ai + 社区技术解读", dark_mode=True)

    # Slide 2: Agenda + method
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, P.light)
    add_title(s, "研究方法与结构导航", "先证据、后结论：文档与源码双轨交叉验证")
    add_card(
        s,
        0.8,
        2.0,
        6.0,
        4.5,
        "分析方法",
        [
            "核心入口：gateway / channels / routing / node-host / media / skills。",
            "证据链：优先本地源码与官方文档；网络资料用于对照视角。",
            "输出目标：提炼“可复用架构机制”，而非功能清单。",
            "强调工程边界：安全、可运维性、演进成本三维评估。",
        ],
        theme="light",
    )
    add_card(
        s,
        7.0,
        2.0,
        5.55,
        4.5,
        "章节结构",
        [
            "A. 架构总览与控制平面",
            "B. 路由/会话/多代理机制",
            "C. 插件化通道与流式回复",
            "D. 安全策略栈与执行审批",
            "E. 模型容灾、媒体管线与运维治理",
            "F. 对 DoWhiz 的可迁移路线图",
        ],
        theme="mint",
    )
    add_footer(s, "目录：18 页，覆盖 11 个核心创新点。")

    # Slide 3: System architecture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, "EAF1FA")
    add_title(s, "系统分层全景（System of Systems）", "Gateway 作为统一中枢，连接渠道、代理、工具与设备节点")

    layers = [
        ("外部消息层", "WhatsApp / Telegram / Discord / Slack / Signal / iMessage / WebChat", "2A9D8F"),
        ("渠道接入层", "channels + channel plugins + outbound adapters", "2E6F95"),
        ("路由会话层", "bindings tiers + session-key algebra + dmScope", "355070"),
        ("Agent 运行层", "pi-agent-core embedded loop + tool streaming", "6D597A"),
        ("工具执行层", "read/write/edit/exec/browser/canvas/nodes", "B56576"),
        ("安全治理层", "pairing/auth/sandbox/policy/elevated/approvals", "E56B6F"),
    ]

    y = 1.95
    for name, desc, c in layers:
        band = s.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(y), Inches(11.2), Inches(0.74))
        band.fill.solid()
        band.fill.fore_color.rgb = rgb(c)
        band.fill.transparency = 6
        band.line.fill.background()

        name_box = s.shapes.add_textbox(Inches(1.3), Inches(y + 0.12), Inches(2.2), Inches(0.36))
        nt = name_box.text_frame
        nt.clear()
        np = nt.paragraphs[0]
        np.text = name
        np.font.name = "PingFang SC"
        np.font.bold = True
        np.font.size = Pt(14)
        np.font.color.rgb = rgb(P.white)

        desc_box = s.shapes.add_textbox(Inches(3.6), Inches(y + 0.12), Inches(8.35), Inches(0.36))
        dt = desc_box.text_frame
        dt.clear()
        dp = dt.paragraphs[0]
        dp.text = desc
        dp.font.name = "Calibri"
        dp.font.size = Pt(13)
        dp.font.color.rgb = rgb("ECF5FF")
        y += 0.85

    add_footer(s, "锚点：docs/concepts/architecture.md + src/gateway/* + src/channels/* + src/routing/*")

    # Slide 4: Innovation 1
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, P.light)
    add_title(s, "创新 1：统一 WebSocket 控制平面", "不是“接口集合”，而是有角色、作用域、挑战握手的协议内核")

    add_card(
        s,
        0.75,
        1.9,
        6.35,
        4.95,
        "关键机制",
        [
            "强类型帧：req/res/event（TypeBox schema）。",
            "握手先发 connect.challenge，再校验 device signature + nonce。",
            "operator 与 node 共用同一协议栈，统一 role/scopes/caps 治理。",
            "side-effect 方法要求 idempotency key，避免重放副作用。",
            "hello-ok 下发 features/snapshot/policy，天然支持多客户端状态对齐。",
        ],
        theme="light",
    )

    flow = s.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(7.35), Inches(2.2), Inches(5.2), Inches(4.2))
    flow.fill.solid()
    flow.fill.fore_color.rgb = rgb("102740")
    flow.line.color.rgb = rgb("294D73")

    f = flow.text_frame
    f.clear()
    fp = f.paragraphs[0]
    fp.text = "握手流程（简化）"
    fp.font.name = "PingFang SC"
    fp.font.bold = True
    fp.font.size = Pt(17)
    fp.font.color.rgb = rgb(P.aqua)

    for line in [
        "① Gateway -> connect.challenge(nonce)",
        "② Client -> req.connect(auth + device + signature)",
        "③ Server: 协议版本/签名/权限/配对状态校验",
        "④ res.hello-ok + deviceToken + policy",
        "⑤ 进入 req/res/event 全双工控制面",
    ]:
        p = f.add_paragraph()
        p.text = line
        p.font.name = "Calibri"
        p.font.size = Pt(13.2)
        p.font.color.rgb = rgb("D8E6F5")

    add_footer(s, "源码：src/gateway/protocol/schema/frames.ts, src/gateway/server/ws-connection/message-handler.ts")

    # Slide 5: Innovation 2 routing & sessions
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, "F5F8FC")
    add_title(s, "创新 2：确定性路由 + 会话键代数", "OpenClaw 把“多渠道上下文一致性”问题工程化为可验证规则")

    add_card(
        s,
        8.7,
        1.92,
        3.85,
        4.95,
        "会话键策略",
        [
            "主键前缀：agent:<agentId>:...",
            "dmScope：main / per-peer / per-channel-peer / per-account-channel-peer。",
            "群组键独立：...:<channel>:group|channel:<id>。",
            "identityLinks 可跨渠道归并同一用户身份。",
        ],
        theme="mint",
    )

    tier_text = [
        "binding.peer（最精确）",
        "binding.peer.parent（线程继承）",
        "binding.guild+roles（Discord 角色路由）",
        "binding.guild",
        "binding.team（Slack）",
        "binding.account",
        "binding.channel",
        "default（兜底代理）",
    ]
    yy = 2.0
    for i, t in enumerate(tier_text, start=1):
        add_tier_row(s, yy, i, t, width=6.6)
        yy += 0.58

    add_footer(s, "源码：src/routing/resolve-route.ts（8级优先级） + src/routing/session-key.ts（键构造）")

    # Slide 6: Innovation 3 pairing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, "F7FAFF")
    add_title(s, "创新 3：双 Pairing 安全模型", "把“消息准入”与“设备准入”拆成两条独立信任链")

    add_card(
        s,
        0.8,
        2.0,
        6.2,
        4.8,
        "DM Pairing（渠道侧）",
        [
            "未知发信人先给 pairing code，不直接触发模型执行。",
            "审批后写入 allowFrom，本地文件持久化。",
            "pairing 请求有 TTL、数量上限与并发锁。",
            "默认策略鼓励 pairing/allowlist，而非开放 DM。",
        ],
        theme="rose",
    )
    add_card(
        s,
        7.15,
        2.0,
        5.4,
        4.8,
        "Device Pairing（节点侧）",
        [
            "设备连接需 device identity + 签名 + nonce 挑战应答。",
            "设备 token 与 role/scopes 绑定，可 rotate/revoke。",
            "approve/reject/remove 全量可审计，广播 pair 事件。",
            "使“节点能力执行”可控且可追溯。",
        ],
        theme="mint",
    )
    add_footer(s, "锚点：docs/channels/pairing.md, src/pairing/pairing-store.ts, src/gateway/server-methods/devices.ts")

    # Slide 7: Innovation 4 plugin architecture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, "EEF4FB")
    add_title(s, "创新 4：通道插件化契约（Channel Plugin Contract）", "把“消息平台差异”收敛为可插拔 adapter 集合")

    matrix = s.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.95), Inches(12.5), Inches(5.0))
    matrix.fill.solid()
    matrix.fill.fore_color.rgb = rgb(P.white)
    matrix.line.color.rgb = rgb("DCE5F0")

    txt = s.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(12.0), Inches(4.4))
    tf = txt.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = "ChannelPlugin 典型适配面"
    p0.font.name = "PingFang SC"
    p0.font.bold = True
    p0.font.size = Pt(20)
    p0.font.color.rgb = rgb(P.ink)

    rows = [
        "config / setup / pairing / security / groups / outbound / status",
        "gateway / auth / commands / threading / streaming / messaging",
        "directory / resolver / actions / heartbeat / agentTools",
        "同一通道生命周期在同一个契约中完成，不再散落在多模块“拼装”。",
        "配合 plugins allow/deny 与 schema 校验，实现扩展能力和治理能力同步增长。",
    ]
    for r in rows:
        pp = tf.add_paragraph()
        pp.text = f"• {r}"
        pp.font.name = "Calibri"
        pp.font.size = Pt(14)
        pp.font.color.rgb = rgb("2C3F58")

    add_footer(s, "源码：src/channels/plugins/types.plugin.ts + types.adapters.ts + docs/tools/plugin.md")

    # Slide 8: Innovation 5 streaming
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, "F4F8FD")
    add_title(s, "创新 5：双层流式回复机制", "Block Streaming 与 Preview Streaming 分层，兼顾可读性与实时反馈")

    add_card(
        s,
        0.8,
        2.0,
        6.1,
        4.8,
        "Block Streaming（真实消息块）",
        [
            "按 text_end 或 message_end 刷出块消息。",
            "min/max 字符阈值 + break preference（段落/换行/句子）。",
            "代码块分割保护：必要时自动闭合与重开 fence。",
            "可配置 coalesce，避免碎片化“连发刷屏”。",
        ],
        theme="light",
    )

    add_card(
        s,
        7.05,
        2.0,
        5.55,
        4.8,
        "Preview Streaming（临时预览）",
        [
            "Telegram/Discord/Slack 支持预览消息编辑。",
            "模式：off / partial / block / progress。",
            "与 block streaming 显式避冲突，防止“双重流式”。",
            "核心是交互体验层，不替代最终消息持久语义。",
        ],
        theme="mint",
    )

    add_footer(s, "文档：docs/concepts/streaming.md")

    # Slide 9: Innovation 6 sandbox/policy/elevated
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, "F2F7FC")
    add_title(s, "创新 6：运行时安全三层栈", "Sandbox（在哪执行）× Tool Policy（能做什么）× Elevated（特权逃逸）")

    tri = s.shapes.add_shape(SHAPE.ISOSCELES_TRIANGLE, Inches(1.0), Inches(1.95), Inches(6.3), Inches(5.1))
    tri.fill.solid()
    tri.fill.fore_color.rgb = rgb("DCEAF8")
    tri.line.color.rgb = rgb("B9D4EF")

    # Layer labels on triangle
    lb1 = s.shapes.add_textbox(Inches(2.25), Inches(2.25), Inches(3.8), Inches(0.5))
    l1 = lb1.text_frame
    l1.clear()
    p = l1.paragraphs[0]
    p.text = "Elevated（仅 exec 特权通道）"
    p.font.name = "PingFang SC"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = rgb("8D3B3B")

    lb2 = s.shapes.add_textbox(Inches(2.05), Inches(3.25), Inches(4.2), Inches(0.5))
    l2 = lb2.text_frame
    l2.clear()
    p = l2.paragraphs[0]
    p.text = "Tool Policy（allow/deny/group）"
    p.font.name = "PingFang SC"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = rgb("2A4A6E")

    lb3 = s.shapes.add_textbox(Inches(2.4), Inches(4.45), Inches(3.6), Inches(0.5))
    l3 = lb3.text_frame
    l3.clear()
    p = l3.paragraphs[0]
    p.text = "Sandbox（host vs docker）"
    p.font.name = "PingFang SC"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = rgb("1E5F55")

    add_card(
        s,
        7.45,
        2.0,
        5.05,
        4.95,
        "工程价值",
        [
            "把“执行位置”和“权限策略”解耦，避免单开关幻觉。",
            "可按 agent/session 维度差异化配置。",
            "对高风险工具默认可收敛到 deny + 审批。",
            "结合 pairing 与 scope，可形成端到端最小权限路径。",
        ],
        theme="light",
    )

    add_footer(s, "文档：docs/gateway/sandboxing.md + docs/gateway/sandbox-vs-tool-policy-vs-elevated.md")

    # Slide 10: Innovation 7 system.run policy engine
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, "F7FAFE")
    add_title(s, "创新 7：system.run 策略引擎（非布尔开关）", "命令执行由 security/ask/allowlist/approval 联合求值")

    stage_w = 2.35
    start_x = 0.95
    y = 2.55
    stages = [
        "命令解析\nargv/shell",
        "allowlist 分析\n匹配与风险判定",
        "审批策略\nask/off/on-miss",
        "执行通道\nhost 或 companion",
        "事件回传\nexec.finished",
    ]
    colors = ["DFF3F2", "E5EEF8", "FDEBE6", "E9F5E5", "EFE8FA"]
    for i, (stage, c) in enumerate(zip(stages, colors)):
        x = start_x + i * (stage_w + 0.2)
        st = s.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(stage_w), Inches(1.6))
        st.fill.solid()
        st.fill.fore_color.rgb = rgb(c)
        st.line.color.rgb = rgb("CFDAE8")
        tf = st.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = stage
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "PingFang SC"
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = rgb(P.ink)
        if i < len(stages) - 1:
            ar = s.shapes.add_shape(SHAPE.CHEVRON, Inches(x + stage_w + 0.03), Inches(y + 0.53), Inches(0.14), Inches(0.5))
            ar.fill.solid()
            ar.fill.fore_color.rgb = rgb(P.gray)
            ar.line.fill.background()

    add_card(
        s,
        0.95,
        4.45,
        11.8,
        2.0,
        "关键点",
        [
            "可识别 shell wrapper（如 bash -c / cmd.exe /c）并阻断绕过式调用。",
            "allow-once / allow-always 与 allowlist 记录联动，形成可追踪演进策略。",
            "配合 exec approval 事件，支持“运行中人类确认”闭环。",
        ],
        theme="light",
    )

    add_footer(s, "源码：src/node-host/exec-policy.ts + src/node-host/invoke-system-run.ts")

    # Slide 11: Innovation 8 model failover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, "F2F7FD")
    add_title(s, "创新 8：模型容灾 = profile 轮转 + fallback 链", "认证状态与模型策略协同，减少单 provider 脆弱性")

    add_card(
        s,
        0.75,
        2.0,
        6.2,
        4.9,
        "Provider 内部",
        [
            "auth profiles（OAuth/API key）按规则轮转。",
            "对 rate-limit/auth/billing 建立 cooldown/disable 状态。",
            "会话维度 profile stickiness，优先稳定缓存与上下文连续性。",
            "失败后优先 provider 内切换，再决定是否跨模型 fallback。",
        ],
        theme="light",
    )

    add_card(
        s,
        7.15,
        2.0,
        5.4,
        4.9,
        "跨模型 fallback",
        [
            "agents.defaults.model.primary + fallbacks 按序尝试。",
            "当 provider 内 profile 已耗尽，切换到下一个模型。",
            "兼顾可靠性与成本，不把失败处理留给外层调用方。",
            "同一机制可支撑多 provider 统一接入。",
        ],
        theme="mint",
    )

    add_footer(s, "文档：docs/concepts/model-failover.md")

    # Slide 12: Innovation 9 compatibility APIs
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, "F6FAFF")
    add_title(s, "创新 9：OpenAI / OpenResponses 兼容网关", "外部协议兼容 + 内部运行时统一，降低生态接入成本")

    api_card = s.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(12.5), Inches(4.9))
    api_card.fill.solid()
    api_card.fill.fore_color.rgb = rgb(P.white)
    api_card.line.color.rgb = rgb("DDE6F3")

    tx = s.shapes.add_textbox(Inches(1.1), Inches(2.25), Inches(11.9), Inches(4.45))
    tf = tx.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = "入口"
    p0.font.name = "PingFang SC"
    p0.font.bold = True
    p0.font.size = Pt(18)
    p0.font.color.rgb = rgb(P.ink)

    for line in [
        "• POST /v1/chat/completions（OpenAI 兼容）",
        "• POST /v1/responses（OpenResponses 兼容，支持 item/tool/image/file 流程）",
        "• 与 Gateway auth、agent routing、sessionKey 规则保持一致",
        "• 请求在内部仍走统一 agentCommand 链路，减少实现分叉",
        "• stream 模式统一映射为 SSE 事件输出",
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "PingFang SC"
        p.font.size = Pt(14)
        p.font.color.rgb = rgb("304862")

    add_footer(s, "源码：src/gateway/openai-http.ts + openresponses-http.ts；文档：docs/gateway/*http-api.md")

    # Slide 13: Innovation 10 skills platform
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, "F3F8FE")
    add_title(s, "创新 10：Skills 平台（分层加载 + 动态刷新）", "把“能力提示词工程”升级为可配置、可门控、可热更新的系统组件")

    stack_x = 1.2
    stack_w = 4.8
    stack_h = 0.72
    stack = [
        ("workspace/skills（最高优先级）", "A7D7F9"),
        ("~/.openclaw/skills（managed）", "CDE9FF"),
        ("bundled skills（内置）", "E2F2FF"),
        ("extraDirs + plugin skills", "F0F8FF"),
    ]
    y0 = 2.2
    for label, c in stack:
        r = s.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(stack_x), Inches(y0), Inches(stack_w), Inches(stack_h))
        r.fill.solid()
        r.fill.fore_color.rgb = rgb(c)
        r.line.color.rgb = rgb("C8DDED")
        tf = r.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.font.name = "PingFang SC"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = rgb(P.ink)
        y0 += 0.83

    add_card(
        s,
        6.35,
        2.0,
        6.2,
        4.9,
        "机制亮点",
        [
            "metadata gate：bins/env/config/os 条件过滤。",
            "session 技能快照复用，降低每轮重建开销。",
            "watcher 监听 SKILL.md 变化，下一轮自动刷新。",
            "可与远端 node 能力联动，动态改变技能可用集合。",
        ],
        theme="light",
    )

    add_footer(s, "文档：docs/tools/skills.md；源码：src/agents/skills/workspace.ts + refresh.ts")

    # Slide 14: Innovation 11 media pipeline
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, "F7FAFD")
    add_title(s, "创新 11：媒体管线的安全与生命周期治理", "从 URL 拉取、类型识别、临时存储到回收，全链路有保护")

    flow_box = s.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(2.0), Inches(12.3), Inches(4.9))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = rgb(P.white)
    flow_box.line.color.rgb = rgb("DDE6F1")

    steps = [
        "URL 输入",
        "SSRF Guard\nDNS/IP/重定向",
        "响应限流\nmaxBytes",
        "MIME 识别\nmagic/header/ext",
        "安全落盘\npath/symlink 防护",
        "TTL 清理\n单次访问回收",
    ]

    x = 1.1
    for i, st in enumerate(steps):
        card = s.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.2), Inches(1.85), Inches(1.45))
        card.fill.solid()
        card.fill.fore_color.rgb = rgb("ECF4FC")
        card.line.color.rgb = rgb("C9DAEC")
        tf = card.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = st
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "PingFang SC"
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = rgb(P.ink)
        if i < len(steps) - 1:
            ar = s.shapes.add_shape(SHAPE.CHEVRON, Inches(x + 1.88), Inches(3.62), Inches(0.16), Inches(0.58))
            ar.fill.solid()
            ar.fill.fore_color.rgb = rgb(P.gray)
            ar.line.fill.background()
        x += 2.0

    note = s.shapes.add_textbox(Inches(1.15), Inches(2.35), Inches(11.7), Inches(0.7))
    nt = note.text_frame
    nt.clear()
    np = nt.paragraphs[0]
    np.text = "重点：OpenClaw 把媒体当作“潜在攻击输入”处理，而非普通附件。"
    np.font.name = "PingFang SC"
    np.font.bold = True
    np.font.size = Pt(16)
    np.font.color.rgb = rgb("274664")

    add_footer(s, "源码：src/media/fetch.ts + store.ts + server.ts")

    # Slide 15: Reliability/Ops
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, "F4F8FD")
    add_title(s, "可靠性与运维：把长期运行作为一等公民", "会话维护、压缩裁剪、重试与幂等共同组成“可持续运行闭环”")

    add_card(
        s,
        0.8,
        2.0,
        4.0,
        4.95,
        "Session 维护",
        [
            "pruneAfter / maxEntries / rotateBytes / maxDiskBytes。",
            "cleanup 支持 warn 与 enforce 模式。",
            "避免会话索引与转录文件无限膨胀。",
        ],
        theme="light",
    )
    add_card(
        s,
        4.95,
        2.0,
        4.0,
        4.95,
        "Context 治理",
        [
            "compaction：持久化摘要进入 JSONL。",
            "session pruning：请求前内存裁剪 toolResult。",
            "二者分工明确，避免上下文失控。",
        ],
        theme="mint",
    )
    add_card(
        s,
        9.1,
        2.0,
        3.45,
        4.95,
        "传输健壮性",
        [
            "provider retry policy（按渠道定制）。",
            "auth/hook 限流与失败回退。",
            "idempotency + dedupe 定时清理。",
        ],
        theme="rose",
    )

    add_footer(s, "文档：docs/concepts/session.md / compaction.md / session-pruning.md / retry.md")

    # Slide 16: Metrics
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, "EEF5FD")
    add_title(s, "工程规模证据（external/openclaw 本地统计）", "不是原型级项目，而是大规模持续演进的生产工程")

    add_big_metric(s, 1.0, 2.2, "3676", "src TypeScript 文件", "1B4965")
    add_big_metric(s, 4.25, 2.2, "1366", "src 测试文件（*test.ts）", "B56576")
    add_big_metric(s, 7.5, 2.2, "652", "docs 文档文件", "2A9D8F")
    add_big_metric(s, 10.75, 2.2, "28MB", "src 体量", "355070")

    band = s.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.1), Inches(12.05), Inches(1.5))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(P.white)
    band.line.color.rgb = rgb("DCE4F0")

    btf = band.text_frame
    btf.clear()
    p = btf.paragraphs[0]
    p.text = "工程启示"
    p.font.name = "PingFang SC"
    p.font.bold = True
    p.font.size = Pt(17)
    p.font.color.rgb = rgb(P.ink)

    for line in [
        "• 复杂度足够高，架构设计重点应放在治理能力与演进稳定性，而非单次功能实现。",
        "• 高测试密度意味着其核心机制（路由、安全、协议）具备较强回归保障基础。",
    ]:
        pp = btf.add_paragraph()
        pp.text = line
        pp.font.name = "PingFang SC"
        pp.font.size = Pt(13.5)
        pp.font.color.rgb = rgb(P.gray)

    add_footer(s, "统计命令：rg --files + wc/du（本地执行）")

    # Slide 17: DoWhiz roadmap
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, "F7FAFE")
    add_title(s, "面向 DoWhiz 的迁移路线图（建议）", "先固化语义层，再引入执行层与安全层，避免“大一统重构风险”")

    roadmap = s.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.0), Inches(12.3), Inches(4.95))
    roadmap.fill.solid()
    roadmap.fill.fore_color.rgb = rgb(P.white)
    roadmap.line.color.rgb = rgb("DCE5F0")

    cols = [
        ("Phase 1（2-4周）", "会话与路由基建", ["定义 session key 规范", "固化 binding tiers", "补回归测试矩阵"]),
        ("Phase 2（4-8周）", "执行与安全治理", ["引入 sandbox/tool policy 分层", "建设 exec 审批链路", "分离 DM 与 device pairing"]),
        ("Phase 3（持续）", "生态与兼容层", ["插件化通道契约", "标准化 HTTP 兼容入口", "skills 与能力目录化"]),
    ]

    cx = 1.2
    for title, subtitle, items in cols:
        card = s.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(2.35), Inches(3.75), Inches(4.25))
        card.fill.solid()
        card.fill.fore_color.rgb = rgb("F3F8FD")
        card.line.color.rgb = rgb("D3E2F1")
        tf = card.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.name = "PingFang SC"
        p0.font.bold = True
        p0.font.size = Pt(15)
        p0.font.color.rgb = rgb(P.ink)

        p1 = tf.add_paragraph()
        p1.text = subtitle
        p1.font.name = "PingFang SC"
        p1.font.size = Pt(12.5)
        p1.font.color.rgb = rgb("3C587A")

        for it in items:
            pp = tf.add_paragraph()
            pp.text = f"• {it}"
            pp.font.name = "PingFang SC"
            pp.font.size = Pt(12.3)
            pp.font.color.rgb = rgb(P.gray)
        cx += 4.05

    add_footer(s, "迁移原则：先语义一致性（路由/会话），再能力扩展（插件/执行）。")

    # Slide 18: Closing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, P.navy, with_deco=True)
    add_title(s, "结论：OpenClaw 的竞争力来自“系统治理能力”", "统一控制面、确定性路由、分层安全策略共同构成其技术护城河", dark_mode=True)

    close = s.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.45), Inches(11.6), Inches(3.9))
    close.fill.solid()
    close.fill.fore_color.rgb = rgb("0F2C4A")
    close.line.color.rgb = rgb("2F567C")

    ctf = close.text_frame
    ctf.clear()
    c0 = ctf.paragraphs[0]
    c0.text = "最终判断"
    c0.font.name = "PingFang SC"
    c0.font.bold = True
    c0.font.size = Pt(22)
    c0.font.color.rgb = rgb(P.aqua)

    points = [
        "OpenClaw 已从“聊天机器人”进化为“个人 AI 网关操作系统”。",
        "其创新重心在架构语义（协议/路由/会话/策略），不是单一模型能力。",
        "对 DoWhiz 的最优学习路径：优先迁移治理框架，再做功能扩展。",
        "这能在复杂度上升时保持可维护、可审计、可持续迭代。",
    ]
    for line in points:
        p = ctf.add_paragraph()
        p.text = f"• {line}"
        p.font.name = "PingFang SC"
        p.font.size = Pt(16)
        p.font.color.rgb = rgb("D5E6F8")

    add_footer(s, "谢谢。附：详细技术文档 openclaw_核心创新技术架构分析.md", dark_mode=True)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_FILE))
    return OUT_FILE


if __name__ == "__main__":
    output = build_ppt()
    print(f"Generated: {output}")
